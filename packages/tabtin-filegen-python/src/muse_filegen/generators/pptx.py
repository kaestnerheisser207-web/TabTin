"""PPTX generator (python-pptx)."""

from __future__ import annotations

from typing import Any, Dict

from pptx import Presentation
from pptx.util import Inches

from muse_filegen.errors import SpecError
from muse_filegen.validate import as_text, optional_text, require_list, require_mapping

_SPEC_HELP = """pptx spec:
{
  "slides": [
    {"title": "Title slide", "subtitle": "Optional subtitle"},
    {"title": "Bullets", "bullets": ["one", "two", "three"]},
    {"title": "Data", "table": {"header": ["c1", "c2"], "rows": [["x", "y"]]}}
  ]
}"""

# python-pptx default template layout indexes.
_LAYOUT_TITLE = 0
_LAYOUT_TITLE_AND_CONTENT = 1
_LAYOUT_TITLE_ONLY = 5


class PptxGenerator:
    file_type = "pptx"
    extensions = (".pptx",)

    def spec_help(self) -> str:
        return _SPEC_HELP

    def generate(self, spec: Dict[str, Any], out_path: str) -> None:
        slides = require_list(spec.get("slides", []), "slides")
        if not slides:
            raise SpecError("pptx spec needs at least one entry in `slides`")

        presentation = Presentation()
        for index, raw_slide in enumerate(slides):
            slide_spec = require_mapping(raw_slide, f"slides[{index}]")
            self._render_slide(presentation, slide_spec, index)

        presentation.save(out_path)

    def _render_slide(
        self, presentation: Presentation, slide_spec: Dict[str, Any], index: int
    ) -> None:
        title = optional_text(slide_spec, "title")
        bullets = slide_spec.get("bullets")
        table = slide_spec.get("table")
        subtitle = slide_spec.get("subtitle")

        if table is not None:
            self._render_table_slide(presentation, title, table, index)
            return

        if bullets is not None:
            self._render_bullets_slide(presentation, title, bullets, index)
            return

        # Title (+ optional subtitle) slide.
        layout = presentation.slide_layouts[_LAYOUT_TITLE]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle is not None and len(slide.placeholders) > 1:
            slide.placeholders[1].text = as_text(subtitle)

    def _render_bullets_slide(
        self, presentation: Presentation, title: str, bullets: Any, index: int
    ) -> None:
        layout = presentation.slide_layouts[_LAYOUT_TITLE_AND_CONTENT]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        items = require_list(bullets, f"slides[{index}].bullets")
        body.clear()
        for item_index, item in enumerate(items):
            paragraph = body.paragraphs[0] if item_index == 0 else body.add_paragraph()
            paragraph.text = as_text(item)

    def _render_table_slide(
        self, presentation: Presentation, title: str, table_spec: Any, index: int
    ) -> None:
        spec = require_mapping(table_spec, f"slides[{index}].table")
        rows = require_list(spec.get("rows", []), f"slides[{index}].table.rows")
        header = spec.get("header")
        header_cells = (
            require_list(header, f"slides[{index}].table.header")
            if header is not None
            else None
        )

        column_count = len(header_cells) if header_cells else _max_row_len(rows)
        if column_count == 0:
            raise SpecError(f"slides[{index}].table has no columns")

        layout = presentation.slide_layouts[_LAYOUT_TITLE_ONLY]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title

        data_rows = [[as_text(c) for c in row] for row in (
            require_list(r, f"slides[{index}].table.rows[{i}]") for i, r in enumerate(rows)
        )]
        all_rows = ([[as_text(c) for c in header_cells]] if header_cells else []) + data_rows
        row_count = max(1, len(all_rows))

        graphic = slide.shapes.add_table(
            row_count,
            column_count,
            Inches(0.5),
            Inches(1.5),
            Inches(9.0),
            Inches(0.4 * row_count),
        )
        grid = graphic.table
        for r in range(row_count):
            source = all_rows[r] if r < len(all_rows) else []
            for c in range(column_count):
                grid.cell(r, c).text = source[c] if c < len(source) else ""


_READ_HELP = """pptx read 输出：
{
  "file_type": "pptx",
  "slides": [{"text": "该页所有文本框文字（换行拼接）"}],
  "text": "所有页用空行拼接"
}"""


class PptxReader:
    file_type = "pptx"
    extensions = (".pptx",)

    def read_help(self) -> str:
        return _READ_HELP

    def read(self, path: str) -> Dict[str, Any]:
        presentation = Presentation(path)
        slides = []
        for slide in presentation.slides:
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            slides.append({"text": "\n".join(texts)})
        return {
            "file_type": "pptx",
            "slides": slides,
            "text": "\n\n".join(slide["text"] for slide in slides),
        }


def _max_row_len(rows: list) -> int:
    return max((len(r) for r in rows if isinstance(r, list)), default=0)
