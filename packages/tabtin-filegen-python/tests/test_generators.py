"""Round-trip tests for each generator, including Chinese (CJK) content."""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from muse_filegen.errors import SpecError
from muse_filegen.registry import get_generator, get_reader


def test_xlsx_round_trip(tmp_path):
    from openpyxl import load_workbook

    out = tmp_path / "t.xlsx"
    get_generator("xlsx").generate(
        {
            "sheets": [
                {
                    "name": "销售数据",
                    "header": ["名称", "数量"],
                    "rows": [["苹果", 3], ["香蕉", 5]],
                    "columns": [{"width": 20}, {"width": 10}],
                }
            ]
        },
        str(out),
    )

    workbook = load_workbook(out)
    sheet = workbook["销售数据"]
    assert sheet["A1"].value == "名称"
    assert sheet["A2"].value == "苹果"
    assert sheet["B2"].value == 3  # numbers stay numeric
    assert sheet.column_dimensions["A"].width == 20


def test_docx_round_trip(tmp_path):
    from docx import Document

    out = tmp_path / "t.docx"
    get_generator("docx").generate(
        {
            "title": "测试文档",
            "blocks": [
                {"type": "heading", "level": 1, "text": "第一章"},
                {"type": "paragraph", "text": "这是中文段落。"},
                {"type": "list", "items": ["项目一", "项目二"]},
                {"type": "table", "header": ["列A", "列B"], "rows": [["甲", "乙"]]},
            ],
        },
        str(out),
    )

    document = Document(out)
    texts = [p.text for p in document.paragraphs]
    assert "测试文档" in texts
    assert "第一章" in texts
    assert "这是中文段落。" in texts
    assert "项目一" in texts
    assert document.tables[0].cell(0, 0).text == "列A"
    assert document.tables[0].cell(1, 1).text == "乙"


def test_pptx_round_trip(tmp_path):
    from pptx import Presentation

    out = tmp_path / "t.pptx"
    get_generator("pptx").generate(
        {
            "slides": [
                {"title": "标题页", "subtitle": "副标题"},
                {"title": "要点", "bullets": ["一", "二", "三"]},
                {"title": "数据", "table": {"header": ["名", "值"], "rows": [["甲", 1]]}},
            ]
        },
        str(out),
    )

    presentation = Presentation(out)
    assert len(presentation.slides) == 3
    all_text = "\n".join(
        shape.text_frame.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "标题页" in all_text
    assert "要点" in all_text
    assert "三" in all_text


def test_pdf_uses_cid_fallback_without_embeddable_font(tmp_path, monkeypatch):
    from muse_filegen.generators import pdf as pdf_module

    monkeypatch.delenv("MUSE_CJK_FONT_PATH", raising=False)
    monkeypatch.delenv("WINDIR", raising=False)
    monkeypatch.setattr(pdf_module, "_SYSTEM_CJK_FONT_PATHS", ())
    pdf_module._ensure_cjk_font.cache_clear()
    try:
        out = tmp_path / "t.pdf"
        get_generator("pdf").generate(
            {
                "title": "中文报告",
                "blocks": [
                    {"type": "heading", "level": 1, "text": "概述"},
                    {"type": "paragraph", "text": "中文不应乱码。"},
                    {"type": "table", "header": ["列A", "列B"], "rows": [["甲", "乙"]]},
                ],
            },
            str(out),
        )

        data = out.read_bytes()
        assert data.startswith(b"%PDF")
        assert b"STSong-Light" in data
    finally:
        pdf_module._ensure_cjk_font.cache_clear()


def test_pdf_rejects_configured_font_without_cjk_glyphs(tmp_path, monkeypatch):
    import reportlab
    from muse_filegen.generators import pdf as pdf_module

    font_path = Path(reportlab.__file__).parent / "fonts" / "Vera.ttf"
    monkeypatch.setenv("MUSE_CJK_FONT_PATH", str(font_path))
    monkeypatch.setattr(pdf_module, "_SYSTEM_CJK_FONT_PATHS", ())
    pdf_module._ensure_cjk_font.cache_clear()
    try:
        out = tmp_path / "fallback.pdf"
        get_generator("pdf").generate(
            {"blocks": [{"type": "paragraph", "text": "中文"}]},
            str(out),
        )

        assert pdf_module._ensure_cjk_font() == "STSong-Light"
        assert b"STSong-Light" in out.read_bytes()
    finally:
        pdf_module._ensure_cjk_font.cache_clear()


def test_pdf_embedded_cjk_text_round_trip_when_system_font_is_available(
    tmp_path, monkeypatch
):
    from pypdf import PdfReader
    from muse_filegen.generators import pdf as pdf_module

    font_path = next(
        (
            Path(candidate)
            for candidate in pdf_module._SYSTEM_CJK_FONT_PATHS
            if Path(candidate).is_file()
        ),
        None,
    )
    if font_path is None:
        pytest.skip("No embeddable system CJK font is available on this host")

    monkeypatch.setenv("MUSE_CJK_FONT_PATH", str(font_path))
    pdf_module._ensure_cjk_font.cache_clear()
    try:
        out = tmp_path / "embedded-cjk.pdf"
        get_generator("pdf").generate(
            {"blocks": [{"type": "paragraph", "text": "中文不应乱码。"}]},
            str(out),
        )

        page = PdfReader(out).pages[0]
        assert "中文不应乱码。" in page.extract_text()
        assert any(_pdf_font_is_embedded(font) for font in page["/Resources"]["/Font"].values())
    finally:
        pdf_module._ensure_cjk_font.cache_clear()


def _pdf_font_is_embedded(font_reference) -> bool:
    font = font_reference.get_object()
    descendants = font.get("/DescendantFonts") or []
    resolved = descendants[0].get_object() if descendants else font
    descriptor = resolved.get("/FontDescriptor")
    if descriptor is None:
        return False
    descriptor = descriptor.get_object()
    return any(key in descriptor for key in ("/FontFile", "/FontFile2", "/FontFile3"))


def test_unknown_block_type_raises(tmp_path):
    out = tmp_path / "bad.docx"
    with pytest.raises(SpecError):
        get_generator("docx").generate(
            {"blocks": [{"type": "nope", "text": "x"}]}, str(out)
        )


def test_xlsx_requires_sheets(tmp_path):
    out = tmp_path / "bad.xlsx"
    with pytest.raises(SpecError):
        get_generator("xlsx").generate({"sheets": []}, str(out))


def test_read_xlsx_round_trip(tmp_path):
    out = tmp_path / "t.xlsx"
    get_generator("xlsx").generate(
        {"sheets": [{"name": "数据", "header": ["名", "值"], "rows": [["甲", 1]]}]},
        str(out),
    )
    content = get_reader("xlsx").read(str(out))
    assert content["file_type"] == "xlsx"
    assert content["sheets"][0]["name"] == "数据"
    assert content["sheets"][0]["rows"] == [["名", "值"], ["甲", 1]]


def test_read_docx_extracts_text(tmp_path):
    out = tmp_path / "t.docx"
    get_generator("docx").generate(
        {"blocks": [{"type": "paragraph", "text": "中文段落"}]}, str(out)
    )
    content = get_reader("docx").read(str(out))
    assert "中文段落" in content["text"]


def test_read_pptx_extracts_text(tmp_path):
    out = tmp_path / "t.pptx"
    get_generator("pptx").generate(
        {"slides": [{"title": "标题页"}, {"title": "要点", "bullets": ["一"]}]}, str(out)
    )
    content = get_reader("pptx").read(str(out))
    assert "标题页" in content["text"]
    assert len(content["slides"]) == 2


def test_read_pdf_extracts_chinese_text(tmp_path):
    out = tmp_path / "t.pdf"
    get_generator("pdf").generate(
        {"blocks": [{"type": "paragraph", "text": "你好世界"}]}, str(out)
    )
    content = get_reader("pdf").read(str(out))
    assert content["file_type"] == "pdf"
    assert "你好世界" in content["text"]


def test_generated_office_files_are_valid_zip(tmp_path):
    for file_type, ext, spec in (
        ("xlsx", ".xlsx", {"sheets": [{"rows": [["a"]]}]}),
        ("docx", ".docx", {"blocks": [{"type": "paragraph", "text": "a"}]}),
        ("pptx", ".pptx", {"slides": [{"title": "a"}]}),
    ):
        out = tmp_path / f"v{ext}"
        get_generator(file_type).generate(spec, str(out))
        assert zipfile.is_zipfile(out)
